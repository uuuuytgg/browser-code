#!/usr/bin/env python3
"""
html2pptx.py - 纯本地 HTML → PPTX 转换器
遵循 pptx skill 的 html2pptx 工作流算法：
  1. 解析 HTML → 拆分幻灯片
  2. 用 python-pptx 构建 PPTX（正确 DEFLATE 压缩）
  3. 保存 .pptx 文件

用法:
  python html2pptx.py input.html output.pptx
  或管道: type file.html | python html2pptx.py output.pptx

依赖: python-pptx (已安装)
"""

import sys, os, re, html as html_mod
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from xml.etree import ElementTree as ET
from io import StringIO

# ====== 模板配色 ======
TEMPLATES = {
    'clean': {
        'name': '干净白底', 'bg': 'FFFFFF', 'bg2': 'F5F5F0',
        'title': '1A1A1A', 'body': '333333',
        'accent': 'DA291C', 'top': 'DA291C',
        'border': 'DDDDDD', 'tblHdr': 'DA291C', 'tblTxt': 'FFFFFF',
    },
    'modern': {
        'name': '现代专业', 'bg': 'FFFFFF', 'bg2': 'F0F4F8',
        'title': '1E293B', 'body': '334155',
        'accent': '3B82F6', 'top': '3B82F6',
        'border': 'CBD5E1', 'tblHdr': '1E293B', 'tblTxt': 'FFFFFF',
    },
    'dark': {
        'name': '深色风格', 'bg': '1E1E2E', 'bg2': '181825',
        'title': 'CDD6F4', 'body': 'A6ADC8',
        'accent': 'F5C2E7', 'top': 'F5C2E7',
        'border': '45475A', 'tblHdr': '45475A', 'tblTxt': 'CDD6F4',
    },
}

TEMPLATE_NAMES = {v['name']: k for k, v in TEMPLATES.items()}

def hex_rgb(h):
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

def emu(inches):
    return int(inches * 914400)

def strip_tags(text):
    return re.sub(r'<[^>]+>', '', text).strip()

def parse_html(html_text, split_mode='h1h2'):
    """从 HTML 中提取幻灯片"""
    # 清理 HTML 实体
    text = html_mod.unescape(html_text)
    # 用正则简单解析（避免依赖 lxml）
    slides = []
    
    # 提取 body 内容
    body_match = re.search(r'<body[^>]*>(.*?)</body>', text, re.DOTALL)
    if body_match:
        body_html = body_match.group(1)
    else:
        body_html = text
    
    # 按模式拆分
    if split_mode == 'section':
        sections = re.findall(r'<section[^>]*>(.*?)</section>', body_html, re.DOTALL)
        if not sections:
            sections = [body_html]
        for sec_html in sections:
            slides.append(_parse_slide_elements(sec_html))
    else:
        split_patterns = {'h1': r'(<h1[^>]*>.*?</h1>)', 'h1h2': r'(<h[12][^>]*>.*?</h[12]>)', 'hr': r'(<h[12][^>]*>.*?</h[12]>|<hr[^>]*>)'}
        pattern = split_patterns.get(split_mode, split_patterns['h1h2'])
        
        parts = re.split(pattern, body_html, flags=re.DOTALL|re.IGNORECASE)
        current = []
        for part in parts:
            part = part.strip()
            if not part: continue
            if re.match(r'<h[12]', part, re.IGNORECASE) or (split_mode=='hr' and re.match(r'<hr', part, re.IGNORECASE)):
                if current:
                    slides.append(_parse_slide_elements('\n'.join(current)))
                current = [part]
            else:
                current.append(part)
        if current:
            slides.append(_parse_slide_elements('\n'.join(current)))
    
    if not slides:
        slides.append(_parse_slide_elements(body_html))
    
    return slides

def _parse_slide_elements(html_text):
    """解析单个幻灯片的元素"""
    slide = {'title': '', 'elements': []}
    
    # 提取标题
    title_match = re.search(r'<h[12][^>]*>(.*?)</h[12]>', html_text, re.DOTALL|re.IGNORECASE)
    if title_match:
        slide['title'] = strip_tags(title_match.group(1))
    
    # 提取段落
    for m in re.finditer(r'<p[^>]*>(.*?)</p>', html_text, re.DOTALL|re.IGNORECASE):
        text = strip_tags(m.group(1))
        if text:
            slide['elements'].append({'type': 'p', 'text': text, 'raw': m.group(1)})
    
    # 提取列表
    for m in re.finditer(r'<(ul|ol)[^>]*>(.*?)</\1>', html_text, re.DOTALL|re.IGNORECASE):
        tag = m.group(1).lower()
        items = [strip_tags(li) for li in re.findall(r'<li[^>]*>(.*?)</li>', m.group(2), re.DOTALL|re.IGNORECASE)]
        slide['elements'].append({'type': tag, 'items': items, 'ordered': tag=='ol'})
    
    # 提取表格
    for m in re.finditer(r'<table[^>]*>(.*?)</table>', html_text, re.DOTALL|re.IGNORECASE):
        rows = []
        for tr in re.findall(r'<tr[^>]*>(.*?)</tr>', m.group(1), re.DOTALL|re.IGNORECASE):
            cells = [strip_tags(td) for td in re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', tr, re.DOTALL|re.IGNORECASE)]
            if cells:
                rows.append(cells)
        slide['elements'].append({'type': 'table', 'rows': rows})
    
    # 提取 blockquote
    for m in re.finditer(r'<blockquote[^>]*>(.*?)</blockquote>', html_text, re.DOTALL|re.IGNORECASE):
        text = strip_tags(m.group(1))
        if text:
            slide['elements'].append({'type': 'quote', 'text': text})
    
    # 提取代码
    for m in re.finditer(r'<(pre|code)[^>]*>(.*?)</\1>', html_text, re.DOTALL|re.IGNORECASE):
        text = html_mod.unescape(m.group(2)).strip()
        if text:
            slide['elements'].append({'type': 'code', 'text': text})
    
    return slide

def build_pptx(slides, template_key='modern', show_num=True):
    """用 python-pptx 构建 PPTX（DEFLATE 压缩，PowerPoint 兼容）"""
    colors = TEMPLATES.get(template_key, TEMPLATES['modern'])
    prs = Presentation()
    prs.slide_width = Emu(emu(13.333))
    prs.slide_height = Emu(emu(7.5))
    
    W = 13.333
    H = 7.5
    total = len(slides)
    
    for i, slide_data in enumerate(slides):
        sd = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        is_first = (i == 0)
        
        # 背景色
        bg = sd.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_rgb(colors['bg'] if not is_first and i%2==0 else (colors['bg2'] if not is_first else colors['bg']))
        
        y = 1.2
        max_y = H - 0.4
        left = 0.6
        content_w = W - 1.2
        
        # 顶部强调线
        _add_rect(sd, 0, 0, W, 0.04, colors['top'])
        
        if is_first:
            # ===== 标题页 =====
            _add_text(sd, slide_data['title'] or '演示文稿', W/2-4, 1.8, 8, 1.5,
                      sz=36, bold=True, color=colors['title'], align='ctr')
            
            sub = next((e for e in slide_data['elements'] if e['type']=='p'), None)
            if sub:
                _add_text(sd, sub['text'], W/2-4, 3.4, 8, 0.7,
                          sz=18, color=colors['body'], align='ctr')
            
            _add_rect(sd, 0, H-0.04, W, 0.04, colors['top'])
        else:
            # ===== 内容页 =====
            # 左侧强调条
            _add_rect(sd, 0.3, 0.3, 0.05, 0.5, colors['accent'])
            
            # 标题
            _add_text(sd, slide_data['title'] or ' ', 0.5, 0.25, W-1, 0.6,
                      sz=22, bold=True, color=colors['title'])
            _add_rect(sd, 0.5, 0.8, 1.5, 0.03, colors['accent'])
            
            # 元素
            for el in slide_data['elements']:
                if y >= max_y: break
                t = el['type']
                if t == 'p':
                    _add_text(sd, el['text'], left, y, content_w, 0.35,
                              sz=13, color=colors['body'])
                    y += 0.35
                elif t in ('ul', 'ol'):
                    prefix = '' if el.get('ordered') else '• '
                    for item in el.get('items', []):
                        if y >= max_y: break
                        _add_text(sd, prefix + item, left+0.2, y, content_w-0.2, 0.3,
                                  sz=12, color=colors['body'])
                        y += 0.28
                elif t == 'table':
                    rows = el.get('rows', [])
                    if rows:
                        _add_table(sd, rows, left, y, content_w, colors)
                        y += min(len(rows) * 0.35, max_y - y) + 0.15
                elif t == 'code':
                    _add_code(sd, el['text'], left, y, content_w, colors)
                    y += 0.55
                elif t == 'quote':
                    _add_quote(sd, el['text'], left, y, content_w, colors)
                    y += 0.5
            
            _add_rect(sd, 0, H-0.04, W, 0.04, colors['top'])
            
            if show_num:
                _add_text(sd, f'{i+1}/{total}', W-1.2, H-0.45, 1, 0.3,
                          sz=8, color=colors['body'], align='r')
    
    return prs

def _add_text(slide, text, x, y, w, h, sz=14, bold=False, color='333333', align='l', italic=False):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h)))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = hex_rgb(color)
    p.font.name = 'Arial'
    if align == 'ctr':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'r':
        p.alignment = PP_ALIGN.RIGHT

def _add_rect(slide, x, y, w, h, color):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(color)
    shape.line.fill.background()

def _add_table(slide, rows, x, y, w, colors):
    """添加表格"""
    cols = max(len(r) for r in rows) if rows else 1
    tbl = slide.shapes.add_table(len(rows), cols,
        Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(min(len(rows)*0.35, 4.0))))
    table = tbl.table
    
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            if ci >= cols: break
            cell = table.cell(ri, ci)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = 'Arial'
                p.font.bold = (ri == 0)
                if ri == 0:
                    p.font.color.rgb = hex_rgb(colors['tblTxt'])
                else:
                    p.font.color.rgb = hex_rgb(colors.get('body', '333333'))
            # 背景
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = hex_rgb(colors['tblHdr'])
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                cell.fill.fore_color.rgb = hex_rgb(colors.get('bg2', 'F5F5F0'))

def _add_code(slide, text, x, y, w, colors):
    """添加代码块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(0.45)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(colors.get('bg2', 'F5F5F0'))
    shape.line.color.rgb = hex_rgb(colors.get('border', 'DDDDDD'))
    shape.line.width = Pt(0.5)
    
    txBox = slide.shapes.add_textbox(Emu(emu(x+0.1)), Emu(emu(y+0.03)),
        Emu(emu(w-0.2)), Emu(emu(0.39)))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = 'Courier New'
    p.font.color.rgb = hex_rgb(colors.get('body', '333333'))

def _add_quote(slide, text, x, y, w, colors):
    """添加引用"""
    _add_rect(slide, x, y, 0.04, 0.4, colors.get('accent', 'DA291C'))
    _add_text(slide, text, x+0.15, y, w-0.15, 0.4,
              sz=12, italic=True, color=colors.get('body', '333333'))

# ====== 主入口 ======
def main():
    if len(sys.argv) < 2:
        print('用法:')
        print('  python html2pptx.py input.html output.pptx')
        print('  python html2pptx.py output.pptx < input.html')
        print('  python html2pptx.py -i input.html -o output.pptx [--template 现代专业|干净白底|深色风格] [--split h1|h1h2|hr|section]')
        return 1
    
    import argparse
    parser = argparse.ArgumentParser(description='HTML → PPTX 纯本地转换器')
    parser.add_argument('-i', '--input', help='输入 HTML 文件')
    parser.add_argument('-o', '--output', default='output.pptx', help='输出 PPTX 文件')
    parser.add_argument('-t', '--template', default='现代专业', choices=['现代专业', '干净白底', '深色风格'])
    parser.add_argument('-s', '--split', default='h1h2', choices=['h1', 'h1h2', 'hr', 'section'])
    args = parser.parse_args()
    
    # 读取输入
    html = ''
    if args.input:
        with open(args.input, 'r', encoding='utf-8') as f:
            html = f.read()
    else:
        html = sys.stdin.read()
    
    if not html.strip():
        print('错误：输入内容为空')
        return 1
    
    # 解析
    template_key = TEMPLATE_NAMES.get(args.template, 'modern')
    slides = parse_html(html, args.split)
    if not slides:
        print('错误：未解析到任何幻灯片内容')
        return 1
    
    print(f'解析到 {len(slides)} 页幻灯片')
    
    # 构建 PPTX
    prs = build_pptx(slides, template_key)
    prs.save(args.output)
    
    size = os.path.getsize(args.output)
    print(f'PPX 已保存: {args.output} ({size/1024:.1f} KB)')
    print(f'模板: {args.template}')
    print(f'幻灯片: {len(slides)} 页')
    print(f'压缩: DEFLATE (PowerPoint 兼容)')
    return 0

if __name__ == '__main__':
    sys.exit(main())
